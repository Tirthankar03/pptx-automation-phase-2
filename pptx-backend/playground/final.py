import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml


def update_title_in_presentation(prs, new_title):
    if not new_title.strip():
        return  # Skip if title is empty

    slide = prs.slides[0]  # First slide
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            count += 1
            if count == 1:  # Update the FIRST text box only
                text_frame = shape.text_frame
                first_paragraph = text_frame.paragraphs[0]
                if first_paragraph.runs:
                    first_run = first_paragraph.runs[0]
                else:
                    first_run = first_paragraph.add_run()

                # Preserve original formatting
                font_name = first_run.font.name
                font_size = first_run.font.size
                font_bold = first_run.font.bold
                font_italic = first_run.font.italic
                font_underline = first_run.font.underline

                # Clear all content
                text_frame.clear()

                # Add new text with preserved formatting
                new_paragraph = text_frame.paragraphs[0]
                new_run = new_paragraph.add_run()
                new_run.text = new_title

                new_run.font.name = font_name
                if font_size:
                    new_run.font.size = font_size
                if font_bold is not None:
                    new_run.font.bold = font_bold
                if font_italic is not None:
                    new_run.font.italic = font_italic
                if font_underline is not None:
                    new_run.font.underline = font_underline
                new_run.font.color.rgb = RGBColor(255, 255, 255)
                return



def clear_table_style(table):
    """Remove default table style to prevent shading."""
    tbl = table._tbl
    tblPr = tbl.get_or_add_tblPr()
    style_id = tblPr.find(qn("a:tableStyleId"))
    if style_id is not None:
        tblPr.remove(style_id)


def set_cell_border(cell, rgb=(0, 0, 0), width_emu="19050"):
    """Apply full borders (all sides) to a cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    r, g, b = rgb
    hex_color = f"{r:02X}{g:02X}{b:02X}"
    line_xml_template = (
        f'<a:ln {nsdecls("a")} w="{width_emu}" cap="flat" cmpd="sng" algn="ctr">'
        f' <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
        f' <a:prstDash val="solid"/>'
        f'</a:ln>'
    )
    for side in ("lnL", "lnR", "lnT", "lnB"):
        existing_ln = tcPr.find(qn(f"a:{side}"))
        if existing_ln is not None:
            tcPr.remove(existing_ln)
        border_element = parse_xml(line_xml_template.replace("a:ln", f"a:{side}"))
        tcPr.append(border_element)


def set_header_border(cell, rgb=(0, 0, 0), width_emu="19050"):
    """Apply borders to header: left, right, bottom only (no top)."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    r, g, b = rgb
    hex_color = f"{r:02X}{g:02X}{b:02X}"
    line_xml_template = (
        f'<a:ln {nsdecls("a")} w="{width_emu}" cap="flat" cmpd="sng" algn="ctr">'
        f' <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
        f' <a:prstDash val="solid"/>'
        f'</a:ln>'
    )
    for side in ("lnL", "lnR", "lnB"):
        existing_ln = tcPr.find(qn(f"a:{side}"))
        if existing_ln is not None:
            tcPr.remove(existing_ln)
        border_element = parse_xml(line_xml_template.replace("a:ln", f"a:{side}"))
        tcPr.append(border_element)


def estimate_lines(text, col_width_emu, font_size_pt=12, avg_char_width_pt=7):
    """Estimate number of lines needed for text wrapping."""
    col_width_pt = col_width_emu / 12700.0
    chars_per_line = max(1, int(col_width_pt / avg_char_width_pt))
    if not text:
        return 1
    words = text.split()
    lines = 0
    current_len = 0
    for word in words:
        word_len = len(word)
        space_needed = 1 if current_len > 0 else 0
        if current_len + space_needed + word_len <= chars_per_line:
            current_len += space_needed + word_len
        else:
            if current_len > 0:
                lines += 1
                current_len = 0
                space_needed = 0
            while word_len > 0:
                chunk_len = min(word_len, chars_per_line - current_len - space_needed)
                current_len += space_needed + chunk_len
                word_len -= chunk_len
                space_needed = 0
                if current_len >= chars_per_line:
                    lines += 1
                    current_len = 0
    if current_len > 0:
        lines += 1
    return max(1, lines)


# JSON data
# json_data = '''
# {
#   "type": "table",
#   "columns": ["Sl no.", "Brief about change", "what is the impact", "Dev effort", "Remarks", "Gone Live/ETA", "Status"],
#   "content": [
#     ["1", "Post sanction edit functionality oisjdgoisjfdg oidsfjgoisdfjg jsdf;gjsdfoi jsdfjgsdfoigjsoifdj gsoidfjosdjf g ", "Reducing manual work and errors", "L", "In UAT testing", "10/11/2025", "action_in_progress"],
#     ["2", "EQUIFAX integration for credit checks", "Improved tracking of CCR reports", "M", "Pending approval", "11/12/2025", "action_yet_to_start"],
#     ["3", "New API for payment processing asdofasdofihasdfi hiashdfiashdf ", "Faster transactions and better security ojasojasd oiajsdfoiajsdd oiajsdfdoiajsddf", "XL", "Development complete", "15/12/2025", "action_over"],
#     ["4", "UI revamp for dashboard aiushdfiausdhf iaushdhfiuashdf iuasdfiuhasdifh iuasdfliuahsdfi hiaushdfiashdfihasdf iuashdfiuashdf ", "Enhanced user experience", "S", "In review", "20/12/2025", "progress_not_as_per_plan"],
#     ["5", "Database optimization  aoshdfouasdhf oahsdfiuahsdf iuuasdffuihasdf", "Improved performance and scalability", "L", "Ongoing", "25/12/2025", "action_in_progress"]
#     ["6", "Database optimization  aoshdfouasdhf oahsdfiuahsdf iuuasdffuihasdf", "Improved performance and scalability", "L", "Ongoing", "25/12/2025", "action_in_progress"]
#   ]
# }
# '''

json_data = '''
{
"type": "table",
"title": "OMPL enhancements NEW 1",
"columns": ["Sl no.", "Brief about change", "what is the impact", "Dev effort", "Remarks", "Gone Live/ETA", "Status"],
"content": [
["1", "96", "84", "2", "60", "10/11/2025", "action_in_progress"],
["1", "96", "84", "2", "60", "10/11/2025", "action_in_progress"],
["1", "96", "84", "2", "60", "10/11/2025", "action_in_progress"],
["1", "96", "84", "2", "60", "10/11/2025", "action_in_progress"],
["1", "96", "84", "2", "60", "10/11/2025", "action_in_progress"],
["1", "96", "84", "2", "60", "10/11/2025", "action_in_progress"],
["1", "96", "84", "2", "60", "10/11/2025", "action_in_progress"],
["1", "96", "84", "2", "60", "10/11/2025", "action_in_progress"],
["1", "96", "84", "2", "60", "10/11/2025", "action_in_progress"],
["1", "96", "84", "2", "60", "10/11/2025", "action_in_progress"],
["1", "96", "84", "2", "60", "10/11/2025", "action_in_progress"]
]
}
'''

data = json.loads(json_data)
columns = data["columns"]
content = data["content"]
title_text = data.get("title", "")

# Max characters validation per column
max_chars = {
    "Sl no.": 1,
    "Brief about change": 96,
    "what is the impact": 84,
    "Dev effort": 2,
    "Remarks": 60,
    "Gone Live/ETA": 10
}

# Validate and truncate if necessary
for row in content:
    for i, cell in enumerate(row):
        col_name = columns[i]
        if col_name in max_chars:
            cell_str = str(cell)
            if len(cell_str) > max_chars[col_name]:
                row[i] = cell_str[:max_chars[col_name]]

# Status mapping and colors
status_map = {
    "action_over": "green",
    "action_in_progress": "blue",
    "progress_not_as_per_plan": "red",
    "action_yet_to_start": "yellow"
}

status_colors = {
    "green": RGBColor(0, 176, 80),    # #00B050
    "blue": RGBColor(0, 112, 192),    # #0070C0
    "red": RGBColor(192, 0, 0),       # #C00000
    "yellow": RGBColor(255, 192, 0)   # #FFC000
}


# Load template
prs = Presentation('powerpoints/template_main_no_table_project_update_footer_new.pptx')

update_title_in_presentation(prs, title_text)

# Table positioning
left = Inches(0.3)
top = Inches(1.2)
width = Inches(12.7)
num_cols = len(columns)

min_row_height = Inches(0.3)

# Column widths
col_widths = [Inches(0.6), Inches(3.2), Inches(2.8), Inches(1.0), Inches(2.0), Inches(1.5), Inches(1.6)]

# Status column index
status_idx = columns.index("Status")

# Alignments
aligns = [PP_ALIGN.CENTER if i == 0 or i == status_idx else PP_ALIGN.LEFT for i in range(num_cols)]

# Split content into chunks of 5 rows per slide
rows_per_slide = 5
chunks = [content[i:i + rows_per_slide] for i in range(0, len(content), rows_per_slide)]

# Prepare slides
slides = [prs.slides[0]]  # Use the first slide from template
for _ in range(len(chunks) - 1):
    slides.append(prs.slides.add_slide(prs.slides[0].slide_layout))

# Process each chunk on its slide
for slide_idx, chunk in enumerate(chunks):
    slide = slides[slide_idx]
    num_rows = len(chunk) + 1
    height = num_rows * min_row_height

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    clear_table_style(table)

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row styling
    for c, head in enumerate(columns):
        cell = table.cell(0, c)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 73, 127)
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.runs[0]
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        set_header_border(cell)

    table.rows[0].height = Inches(0.4)

    # Data rows population
    for r, row_data in enumerate(chunk, start=1):
        for c, text in enumerate(row_data):
            cell = table.cell(r, c)
            if c == status_idx:
                cell.text = ""  # Empty for circle overlay
            else:
                cell.text = text

            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.word_wrap = True

            para = cell.text_frame.paragraphs[0]
            para.alignment = aligns[c]

            if cell.text:
                run = para.runs[0]
                run.font.size = Pt(12)

            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            set_cell_border(cell)

    # Dynamic row heights
    min_height_pt = Inches(0.3).pt
    for r in range(1, num_rows):
        max_lines = 1
        for c in range(num_cols):
            if c == status_idx:
                continue
            cell = table.cell(r, c)
            text = cell.text
            col_width = table.columns[c].width
            lines = estimate_lines(text, col_width)
            max_lines = max(max_lines, lines)

        line_height_pt = 15
        required_height_pt = max_lines * line_height_pt + 20
        table.rows[r].height = Pt(max(required_height_pt, min_height_pt))

    # Circle placement
    circle_diam = Inches(0.25)
    cumulative_y = [top + table.rows[0].height]
    for r in range(1, num_rows):
        cumulative_y.append(cumulative_y[-1] + table.rows[r].height)

    for data_idx, row_data in enumerate(chunk):
        table_row_idx = data_idx + 1

        status_key = row_data[status_idx].lower().replace(" ", "_")
        color_key = status_map.get(status_key, "yellow")

        col_left = left + sum(table.columns[i].width for i in range(status_idx)) + (table.columns[status_idx].width - circle_diam) / 2

        row_top = cumulative_y[data_idx]
        row_center = row_top + table.rows[table_row_idx].height / 2
        circle_top = row_center - circle_diam / 2 + Inches(0.15)

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, col_left, circle_top, circle_diam, circle_diam)
        circle.fill.solid()
        circle.fill.fore_color.rgb = status_colors[color_key]
        circle.line.width = Pt(0)

# Save
output_file = 'final_3.pptx'
if os.path.exists(output_file):
    os.remove(output_file)
prs.save(output_file)
print(f"Created {output_file} - circles now perfectly centered!")