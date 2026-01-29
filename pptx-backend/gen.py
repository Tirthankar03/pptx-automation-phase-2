import os
from pptx import Presentation
from pptx.util import Inches 
from pptx.enum.shapes import MSO_SHAPE

presentation = Presentation(os.path.join("powerpoints", "template_main.pptx"))

def list_text_boxes(presentation, slide_num):
    slide = presentation.slides[slide_num-1]
    text_boxes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            text_boxes.append(shape.text)
    return text_boxes

for idx, text in enumerate(list_text_boxes(presentation, 1), 1):
    print(f"Text Box {idx}: {text}")


def update_text_of_textbox(presentation, slide, text_box_id, new_text):
    slide = presentation.slides[(slide - 1)]
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            count += 1
            if count == text_box_id:
                text_frame = shape.text_frame
                first_paragraph = text_frame.paragraphs[0]
                first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()
                # Preserve formatting of the first run
                font = first_run.font
                font_name = font.name
                font_size = font.size
                font_bold = font.bold
                font_italic = font.italic
                font_underline = font.underline
                # font_color = font.color.rgb
                # Clear existing text and apply new text with preserved formatting
                text_frame.clear()  # Clears all text and formatting
                new_run = text_frame.paragraphs[0].add_run()  # New run in first paragraph
                new_run.text = new_text
                # Reapply formatting
                new_run.font.name = font_name
                new_run.font.size = font_size
                new_run.font.bold = font_bold
                new_run.font.italic = font_italic
                new_run.font.underline = font_underline
                # new_run.font.color.rgb = font_color
                return

uniprot_id = 'iaushdfiuasdfh'
update_text_of_textbox(presentation, 1, 1, uniprot_id)


def add_image_to_slide(slide, image_path, left, top, width, height):
    slide.shapes.add_picture(image_path, left, top, width, height)

# Add images to a specific slide (e.g., slide number 1)
# add_image_to_slide(presentation.slides[1], 'scratch/rna_plot.png', Inches(11), Inches(3), Inches(8) , Inches(6))

left1 = top1 = width1 = height1 = Inches(0.2)

def add_shape_to_slide(slide):
    add_shape1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left1, top1, width1, height1) 
add_shape_to_slide(presentation.slides[0])


# output_path = uniprot_id + "_target_report.pptx"
output_path = "1" + "_target_report.pptx"
presentation.save(output_path)

