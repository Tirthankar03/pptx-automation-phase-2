import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn, nsdecls
import math
import sys



def set_cell_border(cell, rgb=(0, 0, 0), width_emu="19050"):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    r, g, b = rgb
    hex_color = f"{r:02X}{g:02X}{b:02X}"
    line_xml_template = (
        f'<a:ln {nsdecls("a")} w="{width_emu}" cap="flat" cmpd="sng" algn="ctr">'
        f' <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
        f' <a:prstDash val="solid"/>'
        f'</a:ln>'
    )
    for side in ("lnL", "lnR", "lnT", "lnB"):
        existing_ln = tcPr.find(qn(f"a:{side}"))
        if existing_ln is not None:
            tcPr.remove(existing_ln)
        border_element = parse_xml(line_xml_template.replace("a:ln", f"a:{side}"))
        tcPr.append(border_element)


def clear_table_style(table):
    tbl = table._tbl
    tblPr = tbl.get_or_add_tblPr()
    style_id = tblPr.find(qn("a:tableStyleId"))
    if style_id is not None:
        tblPr.remove(style_id)


# def create_dealer_ppt_with_status(filename, title_text, data):
#     prs = Presentation()
#     slide_layout = prs.slide_layouts[5]
#     slide = prs.slides.add_slide(slide_layout)

#     # Title
#     title_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.25), Inches(10.5), Inches(0.75))
#     p = title_box.text_frame.paragraphs[0]
#     p.text = title_text
#     p.alignment = PP_ALIGN.CENTER
#     run = p.runs[0]
#     run.font.size = Pt(28)
#     run.font.bold = True
#     run.font.color.rgb = RGBColor(255, 255, 255)
#     run.font.name = "Arial"
#     title_box.fill.solid()
#     title_box.fill.fore_color.rgb = RGBColor(0, 100, 0)

#     headings = [
#         "SNo", "Change", "Brief about change", "What is impact",
#         "Dev effort", "Gone Live Date", "Remarks", "Status"
#     ]
#     status_col_idx = headings.index("Status")
#     brief_col_idx = headings.index("Brief about change")

#     rows, cols = len(data) + 1, len(headings)
#     table_shape = slide.shapes.add_table(rows, cols, Inches(0.25), Inches(1.25), Inches(10.5), Inches(1.5))
#     table = table_shape.table
#     clear_table_style(table)

#     # --- FIX: scale col widths so all fit inside 10.5 inches ---
#     col_widths_relative = [0.5, 1.0, 3.5, 1.5, 1.0, 1.25, 1.0, 0.75]
#     total_rel = sum(col_widths_relative)
#     scale_factor = 10.5 / total_rel
#     col_widths = [Inches(w * scale_factor) for w in col_widths_relative]
#     for i, width in enumerate(col_widths):
#         table.columns[i].width = width

#     # Header row
#     for col_idx, heading_text in enumerate(headings):
#         cell = table.cell(0, col_idx)
#         cell.text = heading_text
#         cell.fill.solid()
#         cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
#         cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
#         para = cell.text_frame.paragraphs[0]
#         para.alignment = PP_ALIGN.CENTER
#         run = para.runs[0]
#         run.font.size = Pt(14)
#         run.font.bold = True
#         run.font.color.rgb = RGBColor(255, 255, 255)
#         run.font.name = "Arial"
#         set_cell_border(cell, rgb=(0, 0, 0))

#     # Alignment rules
#     alignment_rules = [
#         PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT,
#         PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.CENTER
#     ]

#     # Fill rows
#     for r_idx, row_data in enumerate(data, start=1):
#         for c_idx, cell_text in enumerate(row_data):
#             cell = table.cell(r_idx, c_idx)
#             cell.fill.solid()
#             cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
#             cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
#             if c_idx != status_col_idx:
#                 cell.text = str(cell_text)
#                 para = cell.text_frame.paragraphs[0]
#                 para.alignment = alignment_rules[c_idx]
#                 run = para.runs[0]
#                 run.font.size = Pt(12)
#                 run.font.color.rgb = RGBColor(0, 0, 0)
#                 run.font.name = "Arial"
#             set_cell_border(cell, rgb=(0, 0, 0))

#     # Row height adjustment
#     for r_idx, row_data in enumerate(data, start=1):
#         brief_text = str(row_data[brief_col_idx])
#         chars_per_line_estimate = 70
#         num_lines = math.ceil(len(brief_text) / chars_per_line_estimate)
#         base_height = Inches(0.4)
#         height_per_line = Inches(0.25)
#         calculated_height = base_height + (num_lines * height_per_line)
#         table.rows[r_idx].height = max(Inches(0.6), calculated_height)

#     # Status circles
#     table_left = table_shape.left
#     table_top = table_shape.top
#     final_col_widths = [c.width for c in table.columns]
#     final_row_heights = [r.height for r in table.rows]
#     status_col_left = table_left + sum(final_col_widths[:status_col_idx])
#     current_row_top = table_top + final_row_heights[0]

#     for r_idx, row_data in enumerate(data):
#         status_val = row_data[status_col_idx].lower()
#         cell_width = final_col_widths[status_col_idx]
#         cell_height = final_row_heights[r_idx + 1]
#         circle_diameter = Inches(0.25)
#         circle_left = status_col_left + (cell_width - circle_diameter) / 2
#         circle_top = current_row_top + (cell_height - circle_diameter) / 2

#         if status_val in ("green", "yellow", "red"):
#             circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, circle_left, circle_top, circle_diameter, circle_diameter)
#             circle.fill.solid()
#             if status_val == "green":
#                 circle.fill.fore_color.rgb = RGBColor(0, 176, 80)
#             elif status_val == "yellow":
#                 circle.fill.fore_color.rgb = RGBColor(255, 192, 0)
#             elif status_val == "red":
#                 circle.fill.fore_color.rgb = RGBColor(255, 0, 0)
#             circle.line.width = Pt(1.5)
#             circle.line.color.rgb = RGBColor(0, 0, 0)
#         current_row_top += cell_height

#     # Ensure text wraps properly
#     for row in table.rows:
#         for cell in row.cells:
#             cell.text_frame.word_wrap = True
#             cell.text_frame.auto_size = False

#     prs.save(filename)





def create_dealer_ppt_with_status(filename, title_text, data):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.25), Inches(10.5), Inches(0.75))
    p = title_box.text_frame.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.name = "Arial"
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(0, 100, 0)

    headings = [
        "SNo", "Change", "Brief about change", "What is impact",
        "Dev effort", "Gone Live Date", "Remarks", "Status"
    ]
    status_col_idx = headings.index("Status")
    brief_col_idx = headings.index("Brief about change")

    rows, cols = len(data) + 1, len(headings)
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.25), Inches(1.25), Inches(10.5), Inches(1.5))
    table = table_shape.table
    clear_table_style(table)

    # Scale column widths to 10.5 inches
    col_widths_relative = [0.5, 1.0, 3.5, 1.5, 1.0, 1.25, 1.0, 0.75]
    total_rel = sum(col_widths_relative)
    scale_factor = 10.5 / total_rel
    col_widths = [Inches(w * scale_factor) for w in col_widths_relative]
    for i, width in enumerate(col_widths):
        table.columns[i].width = width

    # Header row
    for col_idx, heading_text in enumerate(headings):
        cell = table.cell(0, col_idx)
        cell.text = heading_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.runs[0]
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.name = "Arial"
        set_cell_border(cell, rgb=(0, 0, 0))

    # Alignment rules
    alignment_rules = [
        PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT,
        PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.CENTER
    ]

    # Fill rows
    for r_idx, row_data in enumerate(data, start=1):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            if c_idx != status_col_idx:
                cell.text = str(cell_text)
                para = cell.text_frame.paragraphs[0]
                para.alignment = alignment_rules[c_idx]
                run = para.runs[0]
                run.font.size = Pt(12)
                run.font.color.rgb = RGBColor(0, 0, 0)
                run.font.name = "Arial"
            set_cell_border(cell, rgb=(0, 0, 0))

    # Row height adjustment
    for r_idx, row_data in enumerate(data, start=1):
        brief_text = str(row_data[brief_col_idx])
        chars_per_line_estimate = 70
        num_lines = math.ceil(len(brief_text) / chars_per_line_estimate)
        base_height = Inches(0.4)
        height_per_line = Inches(0.25)
        calculated_height = base_height + (num_lines * height_per_line)
        table.rows[r_idx].height = max(Inches(0.6), calculated_height)

    # Status circles
    table_left = table_shape.left
    table_top = table_shape.top
    final_col_widths = [c.width for c in table.columns]
    final_row_heights = [r.height for r in table.rows]
    status_col_left = table_left + sum(final_col_widths[:status_col_idx])
    current_row_top = table_top + final_row_heights[0]

    for r_idx, row_data in enumerate(data):
        status_val = str(row_data[status_col_idx]).strip().lower()  # <-- normalized
        cell_width = final_col_widths[status_col_idx]
        cell_height = final_row_heights[r_idx + 1]
        circle_diameter = Inches(0.25)
        circle_left = status_col_left + (cell_width - circle_diameter) / 2
        circle_top = current_row_top + (cell_height - circle_diameter) / 2

        if status_val in ("green", "yellow", "red"):
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, circle_left, circle_top, circle_diameter, circle_diameter)
            circle.fill.solid()
            if status_val == "green":
                circle.fill.fore_color.rgb = RGBColor(0, 176, 80)
            elif status_val == "yellow":
                circle.fill.fore_color.rgb = RGBColor(255, 192, 0)
            else:  # red
                circle.fill.fore_color.rgb = RGBColor(255, 0, 0)
            circle.line.width = Pt(1.5)
            circle.line.color.rgb = RGBColor(0, 0, 0)
        current_row_top += cell_height

    # Ensure text wraps properly
    for row in table.rows:
        for cell in row.cells:
            cell.text_frame.word_wrap = True
            cell.text_frame.auto_size = MSO_AUTO_SIZE.NONE  # <-- fixed
    prs.save(filename)



# --- Tkinter GUI ---
class App(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("PPTX & Excel Generator")
        self.geometry("1200x600")

        self.headings = [
            "SNo", "Change", "Brief about change", "What is impact",
            "Dev effort", "Gone Live Date", "Remarks", "Status"
        ]
        self.data_rows = []

        main_frame = ttk.Frame(self, padding="10")
        main_frame.pack(fill=tk.BOTH, expand=True)

        title_frame = ttk.LabelFrame(main_frame, text="Presentation Title", padding="10")
        title_frame.pack(fill=tk.X, pady=5)
        self.title_entry = ttk.Entry(title_frame, font=("Arial", 12))
        self.title_entry.pack(fill=tk.X, expand=True)
        self.title_entry.insert(0, "Dealer App/Portal Key Enhancements")

        table_container = ttk.LabelFrame(main_frame, text="Data", padding="10")
        table_container.pack(fill=tk.BOTH, expand=True)

        # Notification above the table
        note_label = tk.Label(
            table_container,
            text="Note: 'Brief about change' allows max 350 characters. Other columns allow max 200.",
            fg="red",
            font=("Arial", 9, "italic")
        )
        note_label.pack(anchor="w", pady=3)

        self.canvas = tk.Canvas(table_container)
        scrollbar = ttk.Scrollbar(table_container, orient="vertical", command=self.canvas.yview)
        self.table_frame = ttk.Frame(self.canvas)
        self.table_frame.bind("<Configure>", lambda e: self.canvas.configure(scrollregion=self.canvas.bbox("all")))
        self.canvas.create_window((0, 0), window=self.table_frame, anchor="nw")
        self.canvas.configure(yscrollcommand=scrollbar.set)
        self.canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")

        self.create_table_header()
        self.add_data_row()

        button_frame = ttk.Frame(main_frame, padding="10")
        button_frame.pack(fill=tk.X)
        ttk.Button(button_frame, text="Add Row", command=self.add_data_row).pack(side=tk.LEFT, padx=5)
        ttk.Button(button_frame, text="Remove Last Row", command=self.remove_last_row).pack(side=tk.LEFT, padx=5)
        ttk.Button(button_frame, text="Generate PPTX", command=self.generate_pptx).pack(side=tk.RIGHT, padx=5)
        ttk.Button(button_frame, text="Export to Excel", command=self.export_excel).pack(side=tk.RIGHT, padx=5)

    def validate_entry_length(self, max_chars):
        def validator(P):
            return len(P) <= max_chars
        return (self.register(validator), '%P')

    def create_table_header(self):
        for col_idx, heading in enumerate(self.headings):
            label = ttk.Label(self.table_frame, text=heading, font=("Arial", 10, "bold"), relief=tk.RIDGE, padding=5)
            label.grid(row=0, column=col_idx, sticky="nsew")

    def add_data_row(self):
        row_widgets = []
        row_num = len(self.data_rows) + 1
        for col_idx, heading in enumerate(self.headings):
            max_chars = 350 if heading == "Brief about change" else 200
            entry = ttk.Entry(self.table_frame, font=("Arial", 10))
            entry.grid(row=row_num, column=col_idx, sticky="nsew")
            if heading == "Brief about change":
                def on_validate(P, widget=entry, max_len=max_chars):
                    if len(P) >= max_len:
                        widget.configure(foreground="red")
                    else:
                        widget.configure(foreground="black")
                    return len(P) <= max_len
                vcmd = (self.register(on_validate), "%P")
                entry.configure(validate="key", validatecommand=vcmd)
            else:
                vcmd = self.validate_entry_length(max_chars)
                entry.configure(validate="key", validatecommand=vcmd)
            row_widgets.append(entry)
        self.data_rows.append(row_widgets)
        self.data_rows[-1][0].insert(0, str(row_num))

    def remove_last_row(self):
        if len(self.data_rows) > 1:
            last_row = self.data_rows.pop()
            for widget in last_row:
                widget.destroy()

    def get_data_from_grid(self):
        data = []
        for row_widgets in self.data_rows:
            row_data = [widget.get() for widget in row_widgets]
            if any(row_data):
                data.append(row_data)
        return data

    def generate_pptx(self):
        title = self.title_entry.get()
        data = self.get_data_from_grid()
        if not title:
            messagebox.showerror("Error", "Please enter a presentation title.")
            return
        if not data:
            messagebox.showerror("Error", "Please enter some data.")
            return
        filename = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint Presentations", "*.pptx"), ("All Files", "*.*")]
        )
        if not filename:
            return
        try:
            create_dealer_ppt_with_status(filename, title, data)
            messagebox.showinfo("Success", f"Presentation saved successfully as\n{filename}")
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred 1:\n{e}")

    def export_excel(self):
        data = self.get_data_from_grid()
        if not data:
            messagebox.showerror("Error", "No data to export.")
            return
        filename = filedialog.asksaveasfilename(
            defaultextension=".xlsx",
            filetypes=[("Excel Files", "*.xlsx"), ("All Files", "*.*")]
        )
        if not filename:
            return
        try:
            df = pd.DataFrame(data, columns=self.headings)
            df.to_excel(filename, index=False)
            messagebox.showinfo("Success", f"Excel file saved successfully as\n{filename}")
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred:\n{e}")


if __name__ == "__main__":
    try:
        import pandas
    except ImportError:
        messagebox.showerror("Dependency Error", "Pandas is not installed.\nInstall with: pip install pandas")
        sys.exit(1)

    app = App()
    app.mainloop()