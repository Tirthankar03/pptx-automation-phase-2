from fastapi import FastAPI, UploadFile, File, HTTPException, Form
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml
import io
import tempfile
from pydantic import BaseModel
from typing import List
import pandas as pd
import tempfile
from openpyxl import Workbook
from openpyxl.styles import Font
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.datavalidation import DataValidation
from datetime import datetime
from pptx.enum.shapes import PP_PLACEHOLDER
from copy import deepcopy


app = FastAPI()

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:8080", "http://localhost:8081"],  # Allow your frontend origin
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class PPTXRequest(BaseModel):
    type: str
    title: str
    columns: List[str]
    content: List[List[str]]

# Constants
max_chars = {
    "Sl no.": 4,
    "Brief about change": 96,
    "what is the impact": 84,
    "Dev effort": 2,
    "Remarks": 60,
    "Gone Live/ETA": 10
}

status_map = {
    "action_over": "green",
    "action_in_progress": "blue",
    "progress_not_as_per_plan": "red",
    "action_yet_to_start": "yellow"
}

status_colors = {
    "green": RGBColor(0, 176, 80),    # #00B050
    "blue": RGBColor(0, 112, 192),    # #0070C0
    "red": RGBColor(192, 0, 0),       # #C00000
    "yellow": RGBColor(255, 192, 0)   # #FFC000
}

col_widths = [Inches(0.6), Inches(3.2), Inches(2.8), Inches(1.0), Inches(2.0), Inches(1.5), Inches(1.6)]
rows_per_slide = 5


def generate_pptx(data: PPTXRequest):
    columns = data.columns
    content = data.content
    title_text = data.title

    # Validate
    for row in content:
        for i, cell in enumerate(row):
            col_name = columns[i]
            if col_name in max_chars:
                cell_str = str(cell)
                if len(cell_str) > max_chars[col_name]:
                    raise HTTPException(status_code=400, detail=f"Cell in column '{col_name}' exceeds max length of {max_chars[col_name]} characters.")

    # Load template
    prs = Presentation('powerpoints/template_main_no_table_project_update_footer_new.pptx')

    # update_title_in_presentation(prs, title_text)

    # Table positioning
    left = Inches(0.3)
    top = Inches(1.2)
    width = Inches(12.7)
    num_cols = len(columns)

    min_row_height = Inches(0.3)

    # Status column index
    status_idx = columns.index("Status")

    # Alignments
    aligns = [PP_ALIGN.CENTER if i == 0 or i == status_idx else PP_ALIGN.LEFT for i in range(num_cols)]

    # Split content into chunks
    chunks = [content[i:i + rows_per_slide] for i in range(0, len(content), rows_per_slide)]

    # Prepare slides
    # slides = [prs.slides[0]]
    # for _ in range(len(chunks) - 1):
    #     slides.append(prs.slides.add_slide(prs.slides[0].slide_layout))
    slides = [prs.slides[0]]
    update_title_on_slide(slides[0], title_text)

    # for _ in range(len(chunks) - 1):
    #     new_slide = prs.slides.add_slide(prs.slides[0].slide_layout)
    #     update_title_on_slide(new_slide, title_text)
    #     slides.append(new_slide)
    for _ in range(len(chunks) - 1):
        new_slide = duplicate_slide(prs, prs.slides[0])
        slides.append(new_slide)
        
    # Process each chunk
    for slide_idx, chunk in enumerate(chunks):
        slide = slides[slide_idx]
        num_rows = len(chunk) + 1
        height = num_rows * min_row_height

        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = table_shape.table

        clear_table_style(table)

        for i, w in enumerate(col_widths):
            table.columns[i].width = w

        # Header
        for c, head in enumerate(columns):
            cell = table.cell(0, c)
            cell.text = head
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(30, 73, 127)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.runs[0]
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            set_header_border(cell)

        table.rows[0].height = Inches(0.4)

        # Data rows
        for r, row_data in enumerate(chunk, start=1):
            for c, text in enumerate(row_data):
                cell = table.cell(r, c)
                if c == status_idx:
                    cell.text = ""
                else:
                    cell.text = text

                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.word_wrap = True

                para = cell.text_frame.paragraphs[0]
                para.alignment = aligns[c]

                if cell.text:
                    run = para.runs[0]
                    run.font.size = Pt(12)

                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                set_cell_border(cell)

        # Dynamic heights
        min_height_pt = Inches(0.3).pt
        for r in range(1, num_rows):
            max_lines = 1
            for c in range(num_cols):
                if c == status_idx:
                    continue
                cell = table.cell(r, c)
                text = cell.text
                col_width = table.columns[c].width
                lines = estimate_lines(text, col_width)
                max_lines = max(max_lines, lines)

            line_height_pt = 15
            required_height_pt = max_lines * line_height_pt + 20
            table.rows[r].height = Pt(max(required_height_pt, min_height_pt))

        # Circles
        circle_diam = Inches(0.25)
        cumulative_y = [top + table.rows[0].height]
        for r in range(1, num_rows):
            cumulative_y.append(cumulative_y[-1] + table.rows[r].height)

        for data_idx, row_data in enumerate(chunk):
            table_row_idx = data_idx + 1

            # status_key = row_data[status_idx].lower().replace(" ", "_")
            # color_key = status_map.get(status_key, "yellow")
            color_key = get_status_color(row_data[status_idx])


            col_left = left + sum(table.columns[i].width for i in range(status_idx)) + (table.columns[status_idx].width - circle_diam) / 2

            row_top = cumulative_y[data_idx]
            row_center = row_top + table.rows[table_row_idx].height / 2
            circle_top = row_center - circle_diam / 2 + Inches(0.15)

            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, col_left, circle_top, circle_diam, circle_diam)
            circle.fill.solid()
            circle.fill.fore_color.rgb = status_colors[color_key]
            circle.line.width = Pt(0)

    return prs


@app.post("/generate-pptx")
def generate_pptx_endpoint(request: PPTXRequest):
    prs = generate_pptx(request)
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
        prs.save(tmp.name)
        return FileResponse(tmp.name, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', filename='generated.pptx')


@app.post("/generate-pptx-from-excel")
async def generate_pptx_from_excel(file: UploadFile = File(...)):
    if not file.filename.endswith('.xlsx'):
        raise HTTPException(status_code=400, detail="File must be an .xlsx file.")
    
    # excel_content = await file.read()
    # df = pd.read_excel(io.BytesIO(excel_content))
    # columns = list(df.columns)
    # content = df.astype(str).values.tolist()
    excel_content = await file.read()

    # Read WITHOUT headers
    df_raw = pd.read_excel(io.BytesIO(excel_content), header=None)

    # Title = merged A1:G1
    title = str(df_raw.iloc[0, 0]).strip()

    # Headers = row 2 (A2:G2)
    columns = df_raw.iloc[1].astype(str).tolist()

    if "Status" not in columns:
        raise HTTPException(
            status_code=400,
            detail=f"'Status' column not found. Found columns: {columns}"
        )

    # Data = rows 3 onwards
    # content = (
    #     df_raw.iloc[2:]
    #     .dropna(how="all")      # remove fully empty rows
    #     .astype(str)
    #     .values
    #     .tolist()
    # )
    content = (
        df_raw.iloc[2:]
        .dropna(how="all")
        .applymap(normalize_cell)
        .values
        .tolist()
    )

    # Assuming type is fixed or not needed; set to a default
    data = PPTXRequest(type="project_update", title=title, columns=columns, content=content)
    
    prs = generate_pptx(data)
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
        prs.save(tmp.name)
        return FileResponse(tmp.name, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', filename='generated.pptx')


@app.get("/download-template")
def download_template():
    wb = Workbook()
    ws = wb.active
    ws.title = "Project Update"

    # === Title (row 1, merged) ===
    ws['A1'] = "Project Update Data"
    ws.merge_cells('A1:G1')
    title_cell = ws['A1']
    title_cell.font = Font(bold=True, size=16)

    # === Headers (row 2) ===
    headers = [
        "Sl no.",
        "Brief about change",
        "What is the impact",           # Capitalized for consistency
        "Dev effort",
        "Remarks",
        "Gone Live/ETA",
        "Status"
    ]

    for col_idx, header in enumerate(headers, start=1):
        cell = ws.cell(row=2, column=col_idx, value=header)
        cell.font = Font(bold=True)

    # Optional: pre-fill Sl no. for first 10 rows
    for r in range(3, 13):
        ws.cell(row=r, column=1, value=r - 2)

    # === Dropdown for Status (column G, starting row 3) ===
    status_options = "Action Over,In Progress,Not as per Plan,Yet to Start"
    dv = DataValidation(
        type="list",
        formula1=f'"{status_options}"',
        allow_blank=True,
        showErrorMessage=True,
        errorTitle="Invalid Status",
        error="Please select a valid status from the dropdown."
    )
    dv.add('G3:G500')   # up to 498 data rows — adjust as needed
    ws.add_data_validation(dv)

    # === Better column widths ===
    col_widths = [10, 40, 35, 14, 30, 18, 20]
    for col_idx, width in enumerate(col_widths, start=1):
        ws.column_dimensions[get_column_letter(col_idx)].width = width

    # Freeze header rows so they stay visible when scrolling
    ws.freeze_panes = "A3"

    # Save to temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
        wb.save(tmp.name)

    return FileResponse(
        tmp.name,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        filename="project-update-template.xlsx"
    )
    
# def update_title_in_presentation(prs, new_title):
#     if not new_title.strip():
#         return  # Skip if title is empty

#     slide = prs.slides[0]  # First slide
#     count = 0
#     for shape in slide.shapes:
#         if shape.has_text_frame and shape.text.strip():
#             count += 1
#             if count == 1:  # Update the FIRST text box only
#                 text_frame = shape.text_frame
#                 first_paragraph = text_frame.paragraphs[0]
#                 if first_paragraph.runs:
#                     first_run = first_paragraph.runs[0]
#                 else:
#                     first_run = first_paragraph.add_run()

#                 # Preserve original formatting
#                 font_name = first_run.font.name
#                 font_size = first_run.font.size
#                 font_bold = first_run.font.bold
#                 font_italic = first_run.font.italic
#                 font_underline = first_run.font.underline

#                 # Clear all content
#                 text_frame.clear()

#                 # Add new text with preserved formatting
#                 new_paragraph = text_frame.paragraphs[0]
#                 new_run = new_paragraph.add_run()
#                 new_run.text = new_title

#                 new_run.font.name = font_name
#                 if font_size:
#                     new_run.font.size = font_size
#                 if font_bold is not None:
#                     new_run.font.bold = font_bold
#                 if font_italic is not None:
#                     new_run.font.italic = font_italic
#                 if font_underline is not None:
#                     new_run.font.underline = font_underline
#                 new_run.font.color.rgb = RGBColor(255, 255, 255)
#                 return

def update_title_on_slide(slide, new_title: str):
    if not new_title.strip():
        return

    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            count += 1
            if count == 1:  # first text box = title
                text_frame = shape.text_frame
                first_paragraph = text_frame.paragraphs[0]

                if first_paragraph.runs:
                    first_run = first_paragraph.runs[0]
                else:
                    first_run = first_paragraph.add_run()

                # Preserve formatting
                font_name = first_run.font.name
                font_size = first_run.font.size
                font_bold = first_run.font.bold
                font_italic = first_run.font.italic
                font_underline = first_run.font.underline

                text_frame.clear()

                new_run = text_frame.paragraphs[0].add_run()
                new_run.text = new_title

                new_run.font.name = font_name
                if font_size:
                    new_run.font.size = font_size
                if font_bold is not None:
                    new_run.font.bold = font_bold
                if font_italic is not None:
                    new_run.font.italic = font_italic
                if font_underline is not None:
                    new_run.font.underline = font_underline

                new_run.font.color.rgb = RGBColor(255, 255, 255)
                return



# def update_title_on_slide(slide, new_title: str):
#     if not new_title.strip():
#         return

#     for shape in slide.shapes:
#         if not shape.is_placeholder:
#             continue

#         if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
#             text_frame = shape.text_frame
#             text_frame.clear()

#             p = text_frame.paragraphs[0]
#             run = p.add_run()
#             run.text = new_title

#             # Preserve template styling automatically
#             return

# def update_title_on_slide(slide, new_title: str):
#     if not new_title.strip():
#         return

#     for shape in slide.shapes:
#         if not shape.is_placeholder:
#             continue

#         if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
#             text_frame = shape.text_frame
#             paragraph = text_frame.paragraphs[0]

#             # ---- Capture style BEFORE clearing ----
#             if paragraph.runs:
#                 ref_run = paragraph.runs[0]
#                 font_name = ref_run.font.name
#                 font_size = ref_run.font.size
#                 font_bold = ref_run.font.bold
#                 font_italic = ref_run.font.italic
#                 font_underline = ref_run.font.underline

#                 # IMPORTANT: preserve entire color object
#                 font_color = ref_run.font.color
#             else:
#                 font_name = font_size = font_bold = font_italic = font_underline = None
#                 font_color = None

#             # ---- Clear & reapply ----
#             text_frame.clear()
#             p = text_frame.paragraphs[0]
#             run = p.add_run()
#             run.text = new_title

#             # ---- Restore styling ----
#             if font_name:
#                 run.font.name = font_name
#             if font_size:
#                 run.font.size = font_size
#             if font_bold is not None:
#                 run.font.bold = font_bold
#             if font_italic is not None:
#                 run.font.italic = font_italic
#             if font_underline is not None:
#                 run.font.underline = font_underline

#             # Restore color SAFELY
#             if font_color:
#                 run.font.color._color = font_color._color

#             return


# def update_title_on_slide(slide, new_title: str):
#     if not new_title.strip():
#         return

#     for shape in slide.shapes:
#         if not shape.is_placeholder:
#             continue

#         if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
#             tf = shape.text_frame
#             p = tf.paragraphs[0]

#             # If a run exists, JUST replace its text
#             if p.runs:
#                 p.runs[0].text = new_title
#             else:
#                 # Extremely rare, but safe fallback
#                 run = p.add_run()
#                 run.text = new_title

#             return


# def update_title_on_slide(slide, new_title: str):
#     if not new_title.strip():
#         return

#     for shape in slide.shapes:
#         if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
#             # THIS is the magic line
#             shape.text = new_title
#             return

def clear_table_style(table):
    """Remove default table style to prevent shading."""
    tbl = table._tbl
    tblPr = tbl.get_or_add_tblPr()
    style_id = tblPr.find(qn("a:tableStyleId"))
    if style_id is not None:
        tblPr.remove(style_id)


def set_cell_border(cell, rgb=(0, 0, 0), width_emu="19050"):
    """Apply full borders (all sides) to a cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    r, g, b = rgb
    hex_color = f"{r:02X}{g:02X}{b:02X}"
    line_xml_template = (
        f'<a:ln {nsdecls("a")} w="{width_emu}" cap="flat" cmpd="sng" algn="ctr">'
        f' <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
        f' <a:prstDash val="solid"/>'
        f'</a:ln>'
    )
    for side in ("lnL", "lnR", "lnT", "lnB"):
        existing_ln = tcPr.find(qn(f"a:{side}"))
        if existing_ln is not None:
            tcPr.remove(existing_ln)
        border_element = parse_xml(line_xml_template.replace("a:ln", f"a:{side}"))
        tcPr.append(border_element)


def set_header_border(cell, rgb=(0, 0, 0), width_emu="19050"):
    """Apply borders to header: left, right, bottom only (no top)."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    r, g, b = rgb
    hex_color = f"{r:02X}{g:02X}{b:02X}"
    line_xml_template = (
        f'<a:ln {nsdecls("a")} w="{width_emu}" cap="flat" cmpd="sng" algn="ctr">'
        f' <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
        f' <a:prstDash val="solid"/>'
        f'</a:ln>'
    )
    for side in ("lnL", "lnR", "lnB"):
        existing_ln = tcPr.find(qn(f"a:{side}"))
        if existing_ln is not None:
            tcPr.remove(existing_ln)
        border_element = parse_xml(line_xml_template.replace("a:ln", f"a:{side}"))
        tcPr.append(border_element)


def estimate_lines(text, col_width_emu, font_size_pt=12, avg_char_width_pt=7):
    """Estimate number of lines needed for text wrapping."""
    col_width_pt = col_width_emu / 12700.0
    chars_per_line = max(1, int(col_width_pt / avg_char_width_pt))
    if not text:
        return 1
    words = text.split()
    lines = 0
    current_len = 0
    for word in words:
        word_len = len(word)
        space_needed = 1 if current_len > 0 else 0
        if current_len + space_needed + word_len <= chars_per_line:
            current_len += space_needed + word_len
        else:
            if current_len > 0:
                lines += 1
                current_len = 0
                space_needed = 0
            while word_len > 0:
                chunk_len = min(word_len, chars_per_line - current_len - space_needed)
                current_len += space_needed + chunk_len
                word_len -= chunk_len
                space_needed = 0
                if current_len >= chars_per_line:
                    lines += 1
                    current_len = 0
    if current_len > 0:
        lines += 1
    return max(1, lines)

def get_status_color(status_str: str) -> str:
    if not status_str:
        return "yellow"
    
    s = str(status_str).strip().lower()
    
    mapping = {
        # Exact dropdown values
        "action over": "green",
        "over": "green",
        "completed": "green",
        
        "in progress": "blue",
        "progress": "blue",
        
        "not as per plan": "red",
        "delayed": "red",
        
        "yet to start": "yellow",
        "pending": "yellow",
        "to do": "yellow",
        "todo": "yellow",
    }
    
    for key, color in mapping.items():
        if key in s:
            return color
    
    return "yellow"  # default


def normalize_cell(cell):
    if isinstance(cell, (pd.Timestamp, datetime)):
        return cell.strftime("%d/%m/%Y")  # exactly 10 chars
    return str(cell).strip()


def duplicate_slide(prs, slide):
    slide_layout = slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Remove default shapes
    for shape in list(new_slide.shapes):
        new_slide.shapes._spTree.remove(shape._element)

    # Clone shapes from source slide
    for shape in slide.shapes:
        new_slide.shapes._spTree.append(deepcopy(shape._element))

    return new_slide